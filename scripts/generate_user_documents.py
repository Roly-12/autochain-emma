from __future__ import annotations

import os
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "Utilisateur"
OUTPUT.mkdir(exist_ok=True)

PPTX_PATH = OUTPUT / "Guide_utilisation_AutoChain_Emma.pptx"
PDF_PATH = OUTPUT / "Manuel_complet_AutoChain_Emma.pdf"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
INDIGO = RGBColor(79, 70, 229)
EMERALD = RGBColor(5, 150, 105)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(241, 245, 249)
MUTED = RGBColor(148, 163, 184)


def add_textbox(slide, x, y, w, h, text, size=22, color=WHITE, bold=False,
                align=PP_ALIGN.LEFT, font="Aptos", margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, section, title, number):
    add_textbox(slide, 0.65, 0.35, 9.8, 0.35, section.upper(), 10, RGBColor(129, 140, 248), True)
    add_textbox(slide, 0.65, 0.72, 11.8, 0.7, title, 28, NAVY, True)
    add_textbox(slide, 12.2, 0.38, 0.45, 0.35, f"{number:02}", 11, INDIGO, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.48), Inches(12), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent=INDIGO):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(y + 0.2), Inches(0.08), Inches(h - 0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + 0.42, y + 0.22, w - 0.62, 0.42, title, 15, NAVY, True)
    box = slide.shapes.add_textbox(Inches(x + 0.42), Inches(y + 0.72), Inches(w - 0.62), Inches(h - 0.92))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(body):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE
        p.space_after = Pt(7)
        p.level = 0
        p.text = "• " + p.text


def add_footer(slide):
    add_textbox(
        slide, 0.65, 7.08, 12, 0.22,
        "AutoChain Emma+  |  Guide utilisateur  |  Juillet 2026",
        8, MUTED,
    )


def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = INDIGO
    accent.line.fill.background()
    add_textbox(slide, 0.9, 1.15, 10.5, 0.35, "GUIDE D’UTILISATION", 13, RGBColor(165, 180, 252), True)
    add_textbox(slide, 0.9, 1.65, 11.4, 1.05, "AutoChain Emma+", 42, WHITE, True)
    add_textbox(slide, 0.9, 2.85, 10.5, 0.85, "Gestion de flotte certifiée par la blockchain", 24, RGBColor(203, 213, 225))
    add_textbox(slide, 0.9, 5.7, 10.5, 0.45, "Parcours par rôle · MetaMask · Véhicules · Maintenance · Vente", 14, MUTED)
    add_textbox(slide, 0.9, 6.5, 10.5, 0.3, "Version 1.1 — Déploiement validé en juillet 2026", 10, RGBColor(129, 140, 248), True)

    slides = [
        (
            "Comprendre", "Une application, deux niveaux de confiance",
            [
                ("Données administratives", ["Utilisateurs, photos, échéances et documents privés dans PostgreSQL.", "Accès contrôlé selon le rôle Laravel."], INDIGO),
                ("Preuves blockchain", ["Identifiants techniques, kilométrage, hash de maintenance et transferts.", "Aucune donnée personnelle en clair sur Ethereum."], EMERALD),
                ("Certification réelle", ["Une donnée n’est certifiée qu’après receipt réussi.", "Le serveur vérifie wallet, réseau, contrat, calldata et événement."], RGBColor(234, 88, 12)),
            ],
        ),
        (
            "Premiers pas", "Connexion et navigation",
            [
                ("1. Se connecter", ["Utiliser l’adresse e-mail du compte métier.", "Saisir le code MFA reçu par e-mail : il est obligatoire pour tous."], INDIGO),
                ("2. Lire le tableau de bord", ["Statistiques filtrées selon le rôle.", "Actions rapides et alertes uniquement si autorisées."], EMERALD),
                ("3. Utiliser le menu Plus", ["Documents, carburant, alertes, ventes et utilisateurs.", "Les liens non autorisés sont automatiquement masqués."], RGBColor(14, 116, 144)),
            ],
        ),
        (
            "Rôles", "Qui peut faire quoi ?",
            [
                ("Super Admin", ["Comptes, rôles, véhicule, garage certifié et proposition de vente.", "Wallet identique à l’administrateur du contrat."], INDIGO),
                ("Gestionnaire", ["Données administratives, documents, carburant et suivi de flotte.", "Ne signe pas les opérations réservées au contrat admin."], RGBColor(14, 116, 144)),
                ("Chauffeur", ["Voit uniquement le véhicule affecté.", "Signe un kilométrage croissant et saisit ses pleins."], EMERALD),
                ("Garage / Auditeur", ["Garage : maintenance après certification on-chain.", "Auditeur-acheteur : lecture publique et acceptation de vente."], RGBColor(234, 88, 12)),
            ],
        ),
        (
            "Identité Web3", "Lier MetaMask au compte",
            [
                ("Préparer MetaMask", ["Choisir le réseau demandé.", "Sélectionner le compte correspondant au rôle."], INDIGO),
                ("Signer le challenge", ["Profil → Wallet MetaMask → Lier MetaMask.", "Le challenge expire après cinq minutes et ne coûte aucun ETH."], EMERALD),
                ("Contrôle automatique", ["L’adresse récupérée doit correspondre à la signature.", "Un wallet ne peut être lié qu’à un seul compte AutoChain."], RGBColor(14, 116, 144)),
            ],
        ),
        (
            "Super Admin", "Créer et gérer un véhicule",
            [
                ("Créer", ["Véhicules → Ajouter un véhicule.", "Renseigner VIN, immatriculation, fiche administrative et photo."], INDIGO),
                ("Signer", ["MetaMask affiche registerVehicle.", "Après confirmation, VehicleRegistered apparaît dans la timeline."], EMERALD),
                ("Affecter / statut", ["Affecter un chauffeur dont le wallet est vérifié.", "Signer les statuts critiques : maintenance ou panne."], RGBColor(234, 88, 12)),
            ],
        ),
        (
            "Chauffeur", "Kilométrage et carburant",
            [
                ("Kilométrage", ["Ouvrir le véhicule affecté.", "Saisir une valeur strictement supérieure au dernier relevé."], INDIGO),
                ("Certification", ["Signer updateMileage avec le wallet du chauffeur.", "La valeur change seulement après MileageUpdated."], EMERALD),
                ("Plein", ["Carburant → Saisir un plein.", "Date, litres, montant, station et compteur cohérent."], RGBColor(14, 116, 144)),
            ],
        ),
        (
            "Garage", "Enregistrer une maintenance",
            [
                ("Certification préalable", ["Le Super Admin certifie le wallet du garage on-chain.", "Sans cette étape, aucune maintenance ne peut être signée."], RGBColor(234, 88, 12)),
                ("Intervention", ["Maintenance → Nouvelle maintenance.", "Utiliser exactement le dernier kilométrage certifié."], INDIGO),
                ("Preuve", ["Le détail canonique produit un SHA-256 déterministe.", "MaintenanceRecorded confirme le garage, le hash et le kilométrage."], EMERALD),
            ],
        ),
        (
            "Documents", "Documents, IPFS et alertes",
            [
                ("Documents privés", ["Allowlist PDF/JPEG/PNG/WebP.", "SHA-256 vérifié avant chaque téléchargement et accès journalisé."], INDIGO),
                ("Publication IPFS", ["Le document devient public seulement après CID valide et pinning.", "Le lien gateway est affiché dans l’interface."], EMERALD),
                ("Alertes", ["Contrôle technique, assurance, entretien et seuil kilométrique.", "Déduplication automatique et notification e-mail."], RGBColor(234, 88, 12)),
            ],
        ),
        (
            "Vente", "Double validation du transfert",
            [
                ("Proposition Admin", ["Plus → Ventes → Nouvelle vente.", "Choisir le véhicule et l’acheteur dont le wallet est vérifié."], INDIGO),
                ("Acceptation Acheteur", ["L’acheteur voit uniquement ses ventes.", "Seul le wallet exact peut appeler acceptTransfer."], EMERALD),
                ("Finalisation", ["Le véhicule devient sold après TransferAccepted.", "Le chauffeur est retiré et les deux hashes restent auditables."], RGBColor(234, 88, 12)),
            ],
        ),
        (
            "Traçabilité", "Lire les états et la timeline",
            [
                ("pending", ["Demande créée dans Laravel.", "Aucune certification n’est encore annoncée."], RGBColor(100, 116, 139)),
                ("submitted", ["Transaction envoyée à Ethereum.", "Le receipt est encore attendu ou réconcilié en arrière-plan."], RGBColor(234, 88, 12)),
                ("confirmed", ["Receipt réussi et événement attendu trouvé.", "La donnée certifiée et la timeline sont mises à jour."], EMERALD),
                ("failed", ["Réseau, wallet, calldata, revert ou événement invalide.", "La preuve n’est jamais marquée certifiée."], RGBColor(220, 38, 38)),
            ],
        ),
        (
            "Réseaux", "Hardhat local puis Sepolia",
            [
                ("Hardhat", ["RPC 127.0.0.1:8545, chain ID 31337.", "Comptes publics uniquement pour la démonstration locale."], INDIGO),
                ("Sepolia", ["Chain ID 11155111 et compte dédié non public.", "Chaque signataire paie le gas en Sepolia ETH de test."], EMERALD),
                ("Interdictions", ["Ne jamais utiliser une clé Hardhat sur un réseau public.", "Ne jamais transmettre clé privée ou phrase de récupération."], RGBColor(220, 38, 38)),
            ],
        ),
        (
            "Démarrage Sepolia", "Créer son accès et obtenir des Sepolia ETH",
            [
                ("1. Compte AutoChain", ["Connexion → Créer un compte.", "Le nouveau compte est Auditeur ; le Super Admin attribue un autre rôle si nécessaire."], INDIGO),
                ("2. Wallet MetaMask", ["Installer depuis metamask.io et créer un compte de test dédié.", "Conserver la phrase secrète hors ligne ; ne jamais la communiquer."], RGBColor(14, 116, 144)),
                ("3. Réseau Sepolia", ["MetaMask → Réseaux → Afficher les réseaux de test → Sepolia.", "Vérifier le chain ID 11155111."], EMERALD),
                ("4. ETH de test", ["Copier l’adresse publique et utiliser faucets.chain.link/sepolia.", "Vérifier le solde sur sepolia.etherscan.io puis lier le wallet dans AutoChain."], RGBColor(234, 88, 12)),
            ],
        ),
        (
            "Production", "Application hébergée et services externes",
            [
                ("Application", ["https://autochain-emma.onrender.com", "Render Free exécute une image Docker publique construite par GitHub Actions."], INDIGO),
                ("Données et médias", ["PostgreSQL et stockage S3-compatible sur Supabase.", "Le bucket média public sert avatars, véhicules et logo global."], EMERALD),
                ("E-mail et blockchain", ["Brevo envoie MFA et notifications via HTTPS.", "Sepolia chain ID 11155111 ; chaque signataire paie son gas de test."], RGBColor(234, 88, 12)),
                ("Limites gratuites", ["Render peut mettre le service en veille.", "Queue synchrone et scheduler actif uniquement lorsque le service est éveillé."], RGBColor(100, 116, 139)),
            ],
        ),
        (
            "Bonnes pratiques", "Checklist avant de signer",
            [
                ("Compte", ["Vérifier le nom du compte AutoChain et le rôle affiché.", "Vérifier l’adresse MetaMask sélectionnée."], INDIGO),
                ("Réseau", ["Contrôler le chain ID et l’adresse du contrat.", "Lire la fonction demandée dans l’écran de signature."], EMERALD),
                ("Résultat", ["Attendre confirmed.", "Contrôler le hash, le bloc et l’événement dans la timeline."], RGBColor(14, 116, 144)),
            ],
        ),
    ]

    for number, (section, title, cards) in enumerate(slides, start=2):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT
        add_header(slide, section, title, number)
        count = len(cards)
        width = 11.9 / count
        for index, (card_title, body, accent) in enumerate(cards):
            add_card(slide, 0.65 + index * width, 1.82, width - 0.25, 4.85, card_title, body, accent)
        add_footer(slide)

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, 0.9, 1.25, 11.4, 0.4, "ASSISTANCE", 13, RGBColor(165, 180, 252), True)
    add_textbox(slide, 0.9, 1.9, 11.4, 0.8, "Prêt à utiliser AutoChain Emma+", 34, WHITE, True)
    add_textbox(slide, 0.9, 3.0, 10.5, 1.1, "En cas de doute, ne signez pas. Vérifiez le compte, le réseau et la fonction demandée.", 20, RGBColor(203, 213, 225))
    add_textbox(slide, 0.9, 5.4, 8.5, 0.35, "contact@autochain-emma.com", 14, WHITE, True)
    add_textbox(slide, 0.9, 5.85, 8.5, 0.35, "+242 06 878 18 44", 14, WHITE)
    prs.core_properties.title = "Guide d’utilisation AutoChain Emma+"
    prs.core_properties.subject = "Guide utilisateur par rôle et parcours blockchain"
    prs.core_properties.author = "AutoChain Emma+"
    prs.save(PPTX_PATH)


def make_pdf():
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="CoverTitle", parent=styles["Title"], fontName="Helvetica-Bold",
        fontSize=30, leading=36, textColor=colors.HexColor("#0F172A"),
        alignment=TA_CENTER, spaceAfter=18,
    ))
    styles.add(ParagraphStyle(
        name="CoverSub", parent=styles["Normal"], fontName="Helvetica",
        fontSize=14, leading=20, textColor=colors.HexColor("#475569"),
        alignment=TA_CENTER,
    ))
    styles.add(ParagraphStyle(
        name="H1Custom", parent=styles["Heading1"], fontName="Helvetica-Bold",
        fontSize=20, leading=25, textColor=colors.HexColor("#312E81"),
        spaceBefore=10, spaceAfter=12,
    ))
    styles.add(ParagraphStyle(
        name="H2Custom", parent=styles["Heading2"], fontName="Helvetica-Bold",
        fontSize=14, leading=18, textColor=colors.HexColor("#0F172A"),
        spaceBefore=10, spaceAfter=6,
    ))
    styles.add(ParagraphStyle(
        name="BodyCustom", parent=styles["BodyText"], fontName="Helvetica",
        fontSize=9.6, leading=14.2, textColor=colors.HexColor("#334155"),
        alignment=TA_JUSTIFY, spaceAfter=7,
    ))
    styles.add(ParagraphStyle(
        name="BulletCustom", parent=styles["BodyText"], fontName="Helvetica",
        fontSize=9.4, leading=13.5, leftIndent=12, firstLineIndent=-8,
        textColor=colors.HexColor("#334155"), spaceAfter=4,
    ))
    styles.add(ParagraphStyle(
        name="Note", parent=styles["BodyText"], fontName="Helvetica-Bold",
        fontSize=9.3, leading=13, textColor=colors.HexColor("#92400E"),
        backColor=colors.HexColor("#FEF3C7"), borderPadding=8,
        spaceBefore=7, spaceAfter=9,
    ))

    doc = SimpleDocTemplate(
        str(PDF_PATH), pagesize=A4,
        rightMargin=1.6 * cm, leftMargin=1.6 * cm,
        topMargin=1.8 * cm, bottomMargin=1.7 * cm,
        title="Manuel complet AutoChain Emma+",
        author="AutoChain Emma+",
        subject="Manuel fonctionnel, technique et opérationnel",
    )

    def page_decor(canvas, document):
        canvas.saveState()
        if document.page > 1:
            canvas.setStrokeColor(colors.HexColor("#E2E8F0"))
            canvas.line(1.6 * cm, A4[1] - 1.15 * cm, A4[0] - 1.6 * cm, A4[1] - 1.15 * cm)
            canvas.setFont("Helvetica", 8)
            canvas.setFillColor(colors.HexColor("#64748B"))
            canvas.drawString(1.6 * cm, A4[1] - 0.9 * cm, "AutoChain Emma+ — Manuel complet")
            canvas.drawRightString(A4[0] - 1.6 * cm, 0.85 * cm, f"Page {document.page}")
        canvas.restoreState()

    story = [
        Spacer(1, 2.7 * cm),
        Paragraph("AutoChain Emma+", styles["CoverTitle"]),
        Paragraph("Manuel fonctionnel, technique et opérationnel", styles["CoverSub"]),
        Spacer(1, 1.2 * cm),
        Paragraph(
            "Gestion de flotte, identité wallet, certification Ethereum, documents, "
            "alertes, carburant, maintenance et vente à double validation.",
            styles["CoverSub"],
        ),
        Spacer(1, 3.8 * cm),
        Table(
            [
                ["Version", "1.1"],
                ["Date", "Juillet 2026"],
                ["Environnement", "Laravel 13 · Vue 3 · Render · Supabase · Brevo · Sepolia"],
                ["Contact", "contact@autochain-emma.com · +242 06 878 18 44"],
            ],
            colWidths=[4 * cm, 10.5 * cm],
            style=[
                ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#EEF2FF")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#334155")),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ],
        ),
        PageBreak(),
    ]

    def h1(text):
        story.append(Paragraph(text, styles["H1Custom"]))

    def h2(text):
        story.append(Paragraph(text, styles["H2Custom"]))

    def body(text):
        story.append(Paragraph(text, styles["BodyCustom"]))

    def bullets(items):
        for item in items:
            story.append(Paragraph("• " + item, styles["BulletCustom"]))

    def note(text):
        story.append(Paragraph("IMPORTANT — " + text, styles["Note"]))

    h1("Sommaire")
    toc = [
        ["1", "Objet et périmètre"], ["2", "Travaux réalisés et architecture"],
        ["3", "Installation et exploitation"], ["4", "Rôles et autorisations"],
        ["5", "Authentification, MFA et wallet"], ["6", "Gestion des véhicules"],
        ["7", "Kilométrage et carburant"], ["8", "Maintenance certifiée"],
        ["9", "Documents, IPFS et alertes"], ["10", "Vente à double validation"],
        ["11", "Cycle transactionnel et timeline"], ["12", "Sécurité"],
        ["13", "Hardhat et Sepolia"], ["14", "Tests, CI et livraison"],
        ["15", "Dépannage"], ["16", "Checklist d’exploitation"],
    ]
    story.append(Table(
        toc, colWidths=[1.2 * cm, 13.8 * cm],
        style=[
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
            ("FONTSIZE", (0, 0), (-1, -1), 9.5),
            ("TEXTCOLOR", (0, 0), (0, -1), colors.HexColor("#4F46E5")),
            ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
        ],
    ))
    story.append(PageBreak())

    h1("1. Objet et périmètre")
    body(
        "AutoChain Emma+ est une application de gestion de flotte qui sépare les données "
        "administratives des preuves techniques. PostgreSQL conserve les données métier et "
        "Ethereum conserve des empreintes, états et événements vérifiables. Le présent manuel "
        "décrit l’utilisation quotidienne, les contrôles de sécurité, l’exploitation locale et "
        "le déploiement validé sur https://autochain-emma.onrender.com avec le réseau Sepolia."
    )
    bullets([
        "Périmètre métier : véhicules, affectation, kilométrage, maintenance, documents, carburant, alertes et ventes.",
        "Périmètre Web3 : identité wallet, rôles contractuels, receipts, événements et réconciliation.",
        "Principe central : aucune donnée n’est dite certifiée avant la confirmation du receipt et de l’événement attendu.",
    ])

    h1("2. Travaux réalisés et architecture")
    h2("2.1 Socle applicatif")
    bullets([
        "Laravel 13 avec PostgreSQL, Eloquent, policies, middleware, notifications, scheduler et jobs idempotents.",
        "Vue 3, Inertia.js, Tailwind CSS, interface responsive, thèmes clair/sombre et branding administrateur global.",
        "Routes véhicule sécurisées par UUID et archivage par soft delete.",
        "Matrice d’autorisations centralisée pour les cinq rôles métier.",
    ])
    h2("2.2 Socle blockchain")
    bullets([
        "Contrat VehicleRegistry refondu : admin, garages certifiés, chauffeur affecté et propriétaire.",
        "Kilométrage strictement croissant, statuts critiques, maintenance hashée et transfert en deux transactions.",
        "Suppression du transfert direct qui contournait la double validation.",
        "Événements complets indexés dans la timeline ; aucune PII n’est stockée on-chain.",
    ])
    h2("2.3 Couche de confiance")
    body(
        "Le navigateur envoie les transactions via MetaMask. Laravel ne détient aucune clé privée. "
        "Après soumission du hash, le serveur contrôle le chain ID, l’adresse du contrat, le wallet "
        "émetteur, le calldata exact, le statut du receipt et la présence de l’événement attendu."
    )
    h2("2.4 Architecture hébergée validée")
    bullets([
        "Application Docker sur Render Free, publiée depuis l’image ghcr.io/roly-12/autochain-emma:latest.",
        "PostgreSQL persistant via le Session Pooler Supabase avec TLS obligatoire.",
        "Avatars, photos, logo global et documents stockés dans deux buckets Supabase S3-compatible.",
        "MFA et notifications envoyés par l’API HTTPS Brevo, car Render Free bloque les ports SMTP.",
        "Contrat VehicleRegistry déployé sur Sepolia à l’adresse 0xB04b51e7B65684c409ff45d360342f0a82E18ea0.",
    ])

    h1("3. Installation et exploitation")
    h2("3.1 Prérequis")
    bullets([
        "PHP 8.3 avec pdo_pgsql, bcmath et gd ; Composer ; Node.js ; PostgreSQL ; MetaMask.",
        "Pour le local : Hardhat sur 127.0.0.1:8545, chain ID 31337.",
        "Pour les traitements asynchrones : queue worker et scheduler actifs.",
    ])
    h2("3.2 Services à lancer")
    commands = [
        ["Service", "Commande"],
        ["Laravel", "php artisan serve"],
        ["Frontend", "npm run dev"],
        ["Queue", "php artisan queue:work"],
        ["Scheduler", "php artisan schedule:work"],
        ["Hardhat", "cd blockchain && npm run node"],
        ["Déploiement local", "cd blockchain && npm run deploy:localhost"],
    ]
    story.append(Table(
        commands, colWidths=[4 * cm, 11 * cm],
        repeatRows=1,
        style=[
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#312E81")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTNAME", (0, 1), (-1, -1), "Courier"),
            ("FONTSIZE", (0, 0), (-1, -1), 8.5),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ],
    ))
    h2("3.3 Accès à l’environnement hébergé")
    bullets([
        "URL publique : https://autochain-emma.onrender.com.",
        "Le premier accès après une période d’inactivité peut être ralenti par la veille de Render Free.",
        "Toutes les connexions exigent un code MFA envoyé par Brevo à l’adresse du compte.",
        "Dans MetaMask, sélectionner Sepolia et vérifier le chain ID 11155111 avant toute signature.",
        "Les médias sont servis depuis le bucket public Supabase ; les documents privés restent dans leur bucket dédié.",
    ])

    h1("4. Rôles et autorisations")
    role_data = [
        ["Rôle", "Responsabilités principales"],
        ["Super Admin", "Comptes, véhicule signé, affectation, statut, garage certifié, proposition et annulation de vente."],
        ["Gestionnaire", "Administration de flotte, documents, carburant, alertes et suivi hors opérations contractuelles admin."],
        ["Chauffeur", "Accès au véhicule affecté, kilométrage signé et plein carburant."],
        ["Garagiste agréé", "Maintenance signée uniquement après certification on-chain par l’admin."],
        ["Auditeur / Acheteur", "Lecture des preuves publiques et acceptation de la vente qui lui est destinée."],
    ]
    story.append(Table(
        [[Paragraph(str(c), styles["BodyCustom"]) for c in row] for row in role_data],
        colWidths=[3.8 * cm, 11.2 * cm], repeatRows=1,
        style=[
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#312E81")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ],
    ))
    note("Le wallet ne remplace jamais le compte Laravel. Les deux identités et le rôle doivent correspondre.")

    h1("5. Authentification, MFA et wallet")
    h2("5.1 Connexion applicative")
    bullets([
        "Les comptes inactifs sont bloqués globalement.",
        "Le MFA par e-mail est obligatoire pour tous les rôles et utilise un code à six chiffres valable dix minutes.",
        "Brevo transmet les codes via son API HTTPS ; aucun code ne doit être recherché dans les logs en exploitation normale.",
        "Les sessions sont régénérées après authentification et invalidées à la déconnexion.",
    ])
    h2("5.2 Liaison MetaMask")
    bullets([
        "Ouvrir le menu utilisateur puis Wallet MetaMask.",
        "Demander un challenge unique et signer le message sans transaction ni gas.",
        "Laravel récupère cryptographiquement l’adresse, contrôle l’unicité et marque le wallet vérifié.",
        "Avant chaque transaction, l’application contrôle le réseau et le compte MetaMask sélectionné.",
    ])
    note("Ne jamais transmettre une phrase de récupération ou une clé privée. Les comptes Hardhat sont publics et exclusivement locaux.")

    h1("6. Gestion des véhicules")
    h2("6.1 Création")
    bullets([
        "Seul le Super Admin crée un véhicule on-chain.",
        "Le VIN, l’immatriculation, la marque, le modèle, l’année, le carburant et la photo sont enregistrés hors chaîne.",
        "Un UUID technique est haché puis envoyé à registerVehicle.",
        "Après VehicleRegistered, l’état devient confirmed et le kilométrage initial 0 est certifié.",
    ])
    h2("6.2 Affectation et statut")
    bullets([
        "Le chauffeur doit disposer d’un wallet vérifié.",
        "assignDriver est signé par l’administrateur du contrat.",
        "Les changements critiques utilisent updateStatus et sont visibles dans la timeline.",
        "L’archivage Laravel est un soft delete et ne détruit jamais les preuves blockchain.",
    ])

    h1("7. Kilométrage et carburant")
    h2("7.1 Kilométrage")
    bullets([
        "Le véhicule doit déjà être enregistré on-chain.",
        "La nouvelle valeur doit être strictement supérieure à la valeur certifiée.",
        "Sont autorisés : admin, chauffeur affecté et garage certifié selon le contrat.",
        "La base ne change la valeur certifiée qu’après MileageUpdated confirmé.",
    ])
    h2("7.2 Carburant")
    bullets([
        "Le chauffeur ne peut saisir un plein que pour son véhicule affecté.",
        "Le compteur du plein ne peut pas régresser par rapport aux relevés précédents.",
        "La consommation moyenne utilise la distance entre pleins ordonnés et les litres consommés.",
    ])

    h1("8. Maintenance certifiée")
    bullets([
        "Le Super Admin certifie ou révoque le wallet du garage avec setGarageCertification.",
        "Le garage utilise exactement le dernier kilométrage certifié du véhicule.",
        "Laravel construit un JSON canonique contenant identifiants, type, date, kilométrage, détails, pièces et garage.",
        "Le SHA-256 résultant est identique en base et dans recordMaintenance.",
        "La maintenance devient certified uniquement après MaintenanceRecorded.",
    ])

    h1("9. Documents, IPFS et alertes")
    h2("9.1 Documents")
    bullets([
        "Types autorisés : PDF, JPEG, PNG et WebP, avec taille maximale contrôlée.",
        "Le SHA-256 est recalculé avant téléchargement ; un fichier altéré est bloqué.",
        "Chaque téléchargement ou vérification est inscrit dans le journal d’accès.",
        "Les documents privés sont limités aux rôles de gestion.",
    ])
    h2("9.2 IPFS")
    bullets([
        "Un document n’est public qu’après réception d’un CID valide.",
        "Le service appelle explicitement le pinning avant d’enregistrer is_public=true.",
        "Le CID et le lien gateway sont exposés dans l’interface.",
    ])
    h2("9.3 Alertes")
    bullets([
        "Contrôle technique, assurance, entretien par date et seuil kilométrique.",
        "Une empreinte unique évite les doublons, notamment les alertes de vidange.",
        "Les e-mails sont envoyés par la queue aux responsables actifs ayant activé les notifications.",
    ])

    h1("10. Vente à double validation")
    bullets([
        "Le Super Admin choisit un véhicule enregistré et un acheteur au wallet vérifié.",
        "proposeTransfer produit TransferProposed ; la vente passe admin_signed seulement après confirmation.",
        "Seul le compte Laravel acheteur et son wallet exact peuvent appeler acceptTransfer.",
        "Après TransferAccepted, la vente devient completed, le véhicule sold et le chauffeur est retiré.",
        "Une annulation on-chain utilise cancelTransfer et laisse sa preuve transactionnelle.",
    ])

    h1("11. Cycle transactionnel et timeline")
    tx_data = [
        ["État", "Signification", "Effet métier"],
        ["pending", "Demande préparée", "Aucune donnée certifiée modifiée"],
        ["submitted", "Hash reçu, receipt attendu", "Réconciliation automatique ou manuelle"],
        ["confirmed", "Receipt réussi + événement attendu", "État certifié appliqué atomiquement"],
        ["failed", "Revert ou contrôle invalide", "Preuve refusée, nouvel essai traçable"],
    ]
    story.append(Table(
        [[Paragraph(str(c), styles["BodyCustom"]) for c in row] for row in tx_data],
        colWidths=[2.5 * cm, 5.5 * cm, 7 * cm], repeatRows=1,
        style=[
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#312E81")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
        ],
    ))
    body(
        "Les logs du contrat sont indexés par hash de transaction et index de log. La timeline "
        "fusionne ces événements confirmés avec les informations administratives, les documents "
        "et les pleins, sans présenter ces derniers comme preuves blockchain."
    )

    h1("12. Sécurité")
    bullets([
        "Policies et scopes empêchent les accès inter-rôles et filtrent les tableaux de bord.",
        "Le wallet acheteur, le chauffeur affecté et le garage certifié sont vérifiés avant transaction.",
        "Le calldata serveur doit correspondre exactement à la fonction et aux arguments préparés.",
        "Les clés privées ont été retirées de l’environnement Laravel.",
        "Les fichiers sont stockés hors du répertoire public sauf médias explicitement publiés.",
        "Les comptes inactifs, tentatives MFA et téléchargements sont contrôlés.",
    ])
    note("Toute ancienne clé privée placée dans un environnement applicatif doit être considérée compromise et remplacée.")

    h1("13. Hardhat et Sepolia")
    h2("13.1 Validation locale")
    bullets([
        "Déployer VehicleRegistry sur Hardhat et régénérer deployment.json.",
        "Lier des wallets distincts pour Admin, Garage, Chauffeur et Acheteur.",
        "Exécuter : certification garage, véhicule, affectation, kilométrage, maintenance, proposition et acceptation.",
    ])
    h2("13.2 Sepolia")
    bullets([
        "Utiliser un nouveau compte de test privé, jamais une clé Hardhat.",
        "Configurer un RPC Sepolia fiable et approvisionner chaque signataire via un faucet.",
        "Déployer avec chain ID 11155111 puis relier les wallets des comptes Laravel.",
        "Rejouer tous les scénarios et conserver les liens d’explorateur.",
    ])
    h2("13.3 Créer le compte AutoChain")
    bullets([
        "Sur la page de connexion, cliquer sur « Pas de compte ? Créez-en un ».",
        "Saisir nom, adresse e-mail, mot de passe et confirmation, puis valider l’adresse e-mail si demandé.",
        "Tout nouveau compte reçoit le rôle Auditeur / Acheteur. Le Super Admin doit attribuer Chauffeur, Garagiste agréé ou Gestionnaire lorsque le scénario l’exige.",
        "Le changement de rôle applicatif ne certifie pas automatiquement un garage sur Ethereum.",
    ])
    h2("13.4 Créer un wallet MetaMask de test")
    bullets([
        "Installer uniquement l’extension ou l’application officielle depuis https://metamask.io/download/.",
        "Choisir « Créer un nouveau wallet », définir un mot de passe local et noter la phrase secrète hors ligne.",
        "Créer de préférence un compte séparé pour les essais Sepolia ; ne jamais importer une clé Hardhat.",
        "La phrase secrète et la clé privée ne doivent être saisies ni dans AutoChain, ni dans un faucet, ni envoyées par messagerie.",
    ])
    h2("13.5 Afficher Sepolia")
    bullets([
        "Dans MetaMask, ouvrir Réseaux puis activer « Afficher les réseaux de test ».",
        "Sélectionner « Sepolia » et vérifier le chain ID 11155111.",
        "Si une configuration manuelle est indispensable : symbole ETH et explorateur https://sepolia.etherscan.io/ ; utiliser le RPC fourni par l’hébergeur ou le prestataire retenu.",
    ])
    h2("13.6 Obtenir des Sepolia ETH")
    bullets([
        "Copier uniquement l’adresse publique du wallet actuellement sélectionné.",
        "Essayer un fournisseur reconnu : https://faucets.chain.link/sepolia, https://www.alchemy.com/faucets/ethereum-sepolia ou https://cloud.google.com/application/web3/faucet/ethereum/sepolia.",
        "Coller l’adresse publique, effectuer les contrôles anti-abus demandés puis réclamer les jetons.",
        "Les conditions et limites varient selon le faucet ; essayer un autre fournisseur en cas d’inéligibilité.",
        "Vérifier la réception dans MetaMask et sur https://sepolia.etherscan.io/ en recherchant l’adresse publique.",
    ])
    note("Les Sepolia ETH sont gratuits et n’ont aucune valeur monétaire. Toute personne qui propose de les vendre ou demande la phrase secrète tente probablement une fraude.")
    h2("13.7 Lier le wallet et obtenir les autorisations")
    bullets([
        "Dans AutoChain : menu utilisateur → Wallet MetaMask → Lier MetaMask.",
        "Sélectionner Sepolia et le bon compte, puis signer le challenge gratuit ; cette signature ne consomme pas de gas.",
        "Pour un Chauffeur : le Super Admin doit ensuite l’affecter au véhicule.",
        "Pour un Garagiste : le Super Admin doit attribuer le rôle puis certifier son wallet on-chain.",
        "Pour un Acheteur : le Super Admin doit désigner exactement ce compte dans la proposition de vente.",
        "Chaque transaction métier consomme une petite quantité de Sepolia ETH sur le wallet qui signe.",
    ])
    body(
        "Un testeur externe peut utiliser son propre wallet Sepolia et payer son gas, mais il doit "
        "aussi disposer d’un compte AutoChain actif avec le bon rôle. Un garage doit être certifié, "
        "un chauffeur affecté et un acheteur explicitement désigné."
    )

    h1("14. Tests, CI et livraison")
    bullets([
        "Tests Laravel : comptes inactifs, MFA, rôles, documents, carburant, alertes, transactions et réconciliation.",
        "Tests Solidity : 25 scénarios couvrant permissions, statuts, kilométrage, maintenance et vente.",
        "Test d’intégration Laravel vers Hardhat : chain ID, bytecode du contrat et administrateur.",
        "Build Vite, tests PostgreSQL et publication de l’image Docker exécutés dans GitHub Actions.",
        "Render déploie l’image GHCR afin d’éviter les problèmes de clonage Git depuis la plateforme.",
        "La création d’un véhicule a été signée et confirmée sur Sepolia depuis l’application hébergée.",
    ])

    h1("15. Dépannage")
    troubleshooting = [
        ["Symptôme", "Cause probable", "Action"],
        ["Mauvais compte MetaMask", "Compte sélectionné différent du wallet requis", "Changer de compte, recharger puis signer."],
        ["Réseau incorrect", "Chain ID différent", "Sélectionner Hardhat 31337 ou Sepolia 11155111."],
        ["Réservé à l’administrateur", "Transaction signée par un autre wallet", "Utiliser le wallet déployeur du contrat."],
        ["Maintenance refusée", "Garage non certifié ou km incohérent", "Certifier le garage et reprendre le dernier km confirmé."],
        ["403", "Le rôle ne possède pas la permission", "Vérifier le compte connecté et la matrice des rôles."],
        ["Transaction submitted", "Receipt pas encore disponible", "Attendre le worker ou lancer blockchain:reconcile."],
        ["Document 409", "SHA-256 différent", "Restaurer l’original et auditer le stockage."],
        ["Code MFA absent", "Clé Brevo invalide ou expéditeur non vérifié", "Vérifier BREVO_API_KEY et MAIL_FROM_ADDRESS."],
        ["Image absente", "URL publique ou identifiants S3 Supabase incorrects", "Vérifier le Project URL, les clés S3 et le bucket public."],
        ["Logo absent", "Aucun logo Super Admin ou URL média incorrecte", "Profil Admin → téléverser le logo puis contrôler SUPABASE_MEDIA_URL."],
    ]
    story.append(Table(
        [[Paragraph(str(c), styles["BodyCustom"]) for c in row] for row in troubleshooting],
        colWidths=[3.2 * cm, 5.3 * cm, 6.5 * cm], repeatRows=1,
        style=[
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#312E81")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.45, colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ],
    ))

    h1("16. Checklist d’exploitation")
    h2("Quotidien")
    bullets([
        "Vérifier les alertes critiques et les transactions submitted.",
        "Contrôler le scheduler best effort ; la queue est synchrone sur l’offre Render Free.",
        "Examiner les échecs de notification et les journaux d’accès documentaire.",
    ])
    h2("Avant une démonstration")
    bullets([
        "Vérifier le RPC, le chain ID, le contrat et le solde de chaque wallet.",
        "Confirmer les quatre liaisons wallet et la certification du garage.",
        "Préparer un véhicule neuf pour ne pas réutiliser un identifiant on-chain.",
        "Tester la création, l’affectation, le kilométrage, la maintenance et la vente.",
    ])
    h2("Sauvegarde et continuité")
    bullets([
        "Sauvegarder PostgreSQL Supabase et les deux buckets de stockage.",
        "Conserver deployment.json et les hashes de transaction.",
        "Ne jamais sauvegarder les secrets dans le dépôt Git.",
        "Tester périodiquement la restauration et la commande blockchain:reconcile.",
    ])

    story.append(Spacer(1, 1 * cm))
    story.append(KeepTogether([
        Paragraph("Contact et assistance", styles["H1Custom"]),
        Paragraph(
            "E-mail : contact@autochain-emma.com<br/>"
            "Téléphone : +242 06 878 18 44<br/>"
            "Application : https://autochain-emma.onrender.com",
            styles["BodyCustom"],
        ),
    ]))

    doc.build(story, onFirstPage=page_decor, onLaterPages=page_decor)


if __name__ == "__main__":
    make_presentation()
    make_pdf()
    print(PPTX_PATH)
    print(PDF_PATH)
